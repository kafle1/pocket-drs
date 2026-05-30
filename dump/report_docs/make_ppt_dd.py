"""Build the Defend & Detect final-defense presentation.

Same engine, palette, and slide language as make_ppt.py (the PocketDRS deck):
clean and minimal, navy + one accent blue, a thin accent rule under every
title, a section tag for orientation, generous whitespace. Content swapped for
Sabinesh Rajbhandari's "Defend & Detect: AI-Powered Cybersecurity Platform".
Figures are auto-extracted from his report PDF into dd_figures/.
Requires python-pptx + opencv.
"""

from __future__ import annotations

from pathlib import Path

import cv2
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

BASE = Path(__file__).resolve().parent
TEMPLATE = Path("/Users/nirajkafle/Desktop/niraj/dev-projects/pocket-drs/dump/bit_template.pptx")
FIG = BASE / "dd_figures"
OUT = BASE / "defend_detect_presentation.pptx"

# Restrained, modern palette (identical to the PocketDRS deck).
NAVY = RGBColor(0x12, 0x33, 0x5B)
INK = RGBColor(0x2A, 0x2A, 0x2A)
ACCENT = RGBColor(0x2E, 0x6F, 0xB5)
MUTE = RGBColor(0x7C, 0x86, 0x90)
CARD = RGBColor(0xF1, 0xF5, 0xFA)
HAIR = RGBColor(0xD7, 0xE0, 0xEC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation(str(TEMPLATE))
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[10]

MARGIN = Inches(0.5)
TAG_TOP = Inches(0.42)
TAG_H = Inches(0.22)
TITLE_TOP = Inches(0.62)
TITLE_H = Inches(0.52)
RULE_TOP = Inches(1.16)
BODY_TOP = Inches(1.40)
FOOTER_RESERVE = Inches(0.86)
BODY_H = SH - BODY_TOP - FOOTER_RESERVE
CONTENT_W = SW - 2 * MARGIN

_slide_no = 0


def _new_slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    return s


def _para(tf, i, text, size, bold, color, after=0.0, align=PP_ALIGN.LEFT, lh=1.04):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = text
    p.alignment = align
    p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color
    p.space_after = Pt(after); p.space_before = Pt(0); p.line_spacing = lh
    return p


def _box(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    return tf


def _rect(slide, x, y, w, h, fill, line=None, line_w=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = line_w
    try:
        sp.adjustments[0] = 0.08
    except (IndexError, KeyError):
        pass
    return sp


def _add_number(slide) -> None:
    global _slide_no
    _slide_no += 1
    if _slide_no == 1:
        return
    tf = _box(slide, SW - MARGIN - Inches(0.7), SH - Inches(0.4), Inches(0.7), Inches(0.3))
    _para(tf, 0, str(_slide_no), 9, False, MUTE, align=PP_ALIGN.RIGHT)


def _add_header(slide, title: str, tag: str | None) -> None:
    if tag:
        title_low = title.lower()
        tag_words = [w for w in tag.lower().replace("&", " ").split() if len(w) > 3]
        if any(w in title_low for w in tag_words):
            tag = None
    if tag:
        tf = _box(slide, MARGIN, TAG_TOP, CONTENT_W, TAG_H)
        _para(tf, 0, tag.upper(), 11, True, ACCENT)
    tf = _box(slide, MARGIN, TITLE_TOP, CONTENT_W, TITLE_H, anchor=MSO_ANCHOR.MIDDLE)
    _para(tf, 0, title, 23, True, NAVY)
    _rect(slide, MARGIN, RULE_TOP, Inches(1.9), Pt(3), ACCENT, shape=MSO_SHAPE.RECTANGLE)
    _add_number(slide)


def _lead(slide, text: str):
    h = Inches(0.5)
    tf = _box(slide, MARGIN, BODY_TOP, CONTENT_W, h, anchor=MSO_ANCHOR.MIDDLE)
    _para(tf, 0, text, 16, True, NAVY, lh=1.02)
    return BODY_TOP + h + Inches(0.12)


# ----------------------------------------------------------------------------
# Slide templates (verbatim from the PocketDRS deck)
# ----------------------------------------------------------------------------

def content_slide(tag: str, title: str, lead: str, bullets: list[str]) -> None:
    s = _new_slide()
    _add_header(s, title, tag)
    top = _lead(s, lead)
    tf = _box(s, MARGIN, top, CONTENT_W, (BODY_TOP + BODY_H) - top)
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r1 = p.add_run(); r1.text = "▪  "
        r1.font.size = Pt(13); r1.font.color.rgb = ACCENT; r1.font.bold = True
        r2 = p.add_run(); r2.text = text
        r2.font.size = Pt(15); r2.font.color.rgb = INK
        p.space_after = Pt(11); p.line_spacing = 1.06


def pipeline_slide(tag: str, title: str, lead: str, steps: list[tuple[str, str]]) -> None:
    s = _new_slide()
    _add_header(s, title, tag)
    top = _lead(s, lead)
    n = len(steps)
    gap = Inches(0.16)
    card_w = int((CONTENT_W - gap * (n - 1)) / n)
    band_top = top + Inches(0.08)
    band_h = int((BODY_TOP + BODY_H) - band_top)
    card_h = min(band_h, int(Inches(2.15)))
    card_top = band_top + (band_h - card_h) // 2
    badge = Inches(0.42)
    x = MARGIN
    for k, (name, desc) in enumerate(steps):
        _rect(s, x, card_top, card_w, card_h, CARD, line=HAIR, line_w=Pt(0.75))
        bx = x + (card_w - badge) // 2
        _rect(s, bx, card_top + Inches(0.22), badge, badge, ACCENT, shape=MSO_SHAPE.OVAL)
        btf = _box(s, bx, card_top + Inches(0.22), badge, badge, anchor=MSO_ANCHOR.MIDDLE)
        _para(btf, 0, str(k + 1), 16, True, WHITE, align=PP_ALIGN.CENTER)
        ty = card_top + Inches(0.74)
        ttf = _box(s, x + Inches(0.1), ty, card_w - Inches(0.2),
                   (card_top + card_h) - ty - Inches(0.08), anchor=MSO_ANCHOR.MIDDLE)
        _para(ttf, 0, name, 12.5, True, NAVY, after=4, align=PP_ALIGN.CENTER, lh=1.0)
        _para(ttf, 1, desc, 10.5, False, INK, align=PP_ALIGN.CENTER, lh=1.06)
        if k < n - 1:
            cf = _box(s, x + card_w - Inches(0.02), card_top, gap + Inches(0.04), card_h,
                      anchor=MSO_ANCHOR.MIDDLE)
            _para(cf, 0, "›", 20, True, ACCENT, align=PP_ALIGN.CENTER)
        x += card_w + gap


def cards_slide(tag: str, title: str, lead: str, cards: list[tuple[str, str]]) -> None:
    s = _new_slide()
    _add_header(s, title, tag)
    top = _lead(s, lead)
    cols, gap = 2, Inches(0.22)
    rows = (len(cards) + cols - 1) // cols
    cw = int((CONTENT_W - gap * (cols - 1)) / cols)
    grid_top = top + Inches(0.06)
    ch = int(((BODY_TOP + BODY_H) - grid_top - gap * (rows - 1)) / rows)
    for idx, (term, meaning) in enumerate(cards):
        r, c = divmod(idx, cols)
        x = MARGIN + c * (cw + gap)
        y = grid_top + r * (ch + gap)
        _rect(s, x, y, cw, ch, CARD, line=HAIR, line_w=Pt(0.75))
        tf = _box(s, x + Inches(0.22), y, cw - Inches(0.44), ch, anchor=MSO_ANCHOR.MIDDLE)
        _para(tf, 0, term, 14, True, NAVY, after=4, lh=1.0)
        _para(tf, 1, meaning, 12, False, INK, lh=1.07)


def _fit(img: Path, max_w: int, max_h: int) -> tuple[int, int]:
    h, w = cv2.imread(str(img)).shape[:2]
    ar = w / h
    width, height = max_w, int(max_w / ar)
    if height > max_h:
        height, width = max_h, int(max_h * ar)
    return width, height


def image_slide(tag: str, title: str, images: list[Path],
                captions: list[str] | None = None) -> None:
    images = [im for im in images if im.exists()]
    if not images:
        return
    s = _new_slide()
    _add_header(s, title, tag)
    n = len(images)
    gap = Inches(0.3)
    cap_h = Inches(0.3) if captions else Inches(0.0)
    cell_w = int((CONTENT_W - gap * (n - 1)) / n)
    area_h = int(BODY_H - cap_h)
    x = MARGIN
    for k, img in enumerate(images):
        w, h = _fit(img, cell_w, area_h)
        s.shapes.add_picture(str(img), x + (cell_w - w) // 2,
                             BODY_TOP + (area_h - h) // 2, width=w, height=h)
        if captions:
            cf = _box(s, x, BODY_TOP + area_h, cell_w, cap_h)
            _para(cf, 0, captions[k], 11, False, MUTE, align=PP_ALIGN.CENTER)
        x += cell_w + gap


def stat_image_slide(tag: str, title: str, image: Path, lead: str,
                     stats: list[tuple[str, str]], verdict: str | None = None) -> None:
    s = _new_slide()
    _add_header(s, title, tag)
    top = _lead(s, lead)
    area_h = (BODY_TOP + BODY_H) - top
    half = int((CONTENT_W - Inches(0.35)) / 2)
    if image.exists():
        w, h = _fit(image, half, int(area_h))
        s.shapes.add_picture(str(image), MARGIN + (half - w) // 2,
                             top + (int(area_h) - h) // 2, width=w, height=h)
    rx = MARGIN + half + Inches(0.35)
    gap = Inches(0.16)
    pill_h = Inches(0.62) if verdict else Inches(0.0)
    pill_gap = Inches(0.16) if verdict else Inches(0.0)
    grid_h = int(area_h - pill_h - pill_gap)
    rows = (len(stats) + 1) // 2
    cw = int((half - gap) / 2)
    chh = int((grid_h - gap * (rows - 1)) / rows)
    for idx, (val, label) in enumerate(stats):
        r, c = divmod(idx, 2)
        x = rx + c * (cw + gap)
        y = top + r * (chh + gap)
        _rect(s, x, y, cw, chh, CARD, line=HAIR, line_w=Pt(0.75))
        tf = _box(s, x + Inches(0.12), y, cw - Inches(0.24), chh, anchor=MSO_ANCHOR.MIDDLE)
        _para(tf, 0, val, 21, True, NAVY, after=2, lh=0.95)
        _para(tf, 1, label, 10.5, False, MUTE, lh=1.0)
    if verdict:
        py = top + area_h - pill_h
        _rect(s, rx, py, half, pill_h, ACCENT)
        tf = _box(s, rx + Inches(0.2), py, half - Inches(0.4), pill_h, anchor=MSO_ANCHOR.MIDDLE)
        _para(tf, 0, verdict, 14, True, WHITE, align=PP_ALIGN.CENTER, lh=1.0)


def title_slide() -> None:
    s = _new_slide()
    _add_number(s)
    _rect(s, MARGIN, Inches(1.05), Inches(2.2), Pt(4), ACCENT, shape=MSO_SHAPE.RECTANGLE)
    tf = _box(s, MARGIN, Inches(1.2), CONTENT_W, Inches(3.2))
    rows = [
        ("Defend & Detect", 42, True, NAVY, 6),
        ("An AI-Powered Cybersecurity Platform:", 17, False, INK, 0),
        ("A Unified Forensic Workbench for Threat Intelligence", 17, False, INK, 18),
        ("Sabinesh Rajbhandari,  BIT 7th Semester,  ID LC0003001698", 14, False, INK, 3),
        ("Supervisor: Mr. Saishab Bhattarai", 13, False, INK, 3),
        ("Phoenix College of Management,  Lincoln University College", 13, False, ACCENT, 0),
    ]
    for i, (t, sz, b, c, after) in enumerate(rows):
        _para(tf, i, t, sz, b, c, after=after, lh=1.02)


def agenda_slide() -> None:
    s = _new_slide()
    _add_header(s, "Agenda", None)
    items = [
        "Introduction and the Problem",
        "Objectives and Requirements",
        "System Analysis and Design",
        "The Analysis Pipeline",
        "Key Terms",
        "Implementation and Tools",
        "Results and Testing",
        "Limitations and Conclusion",
    ]
    tf = _box(s, MARGIN, BODY_TOP, CONTENT_W, BODY_H, anchor=MSO_ANCHOR.MIDDLE)
    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run(); run.text = f"{i + 1:02d}    "
        run.font.size = Pt(16); run.font.bold = True; run.font.color.rgb = ACCENT
        run2 = p.add_run(); run2.text = text
        run2.font.size = Pt(16); run2.font.color.rgb = INK
        p.space_after = Pt(8); p.line_spacing = 1.05


def closing_slide() -> None:
    s = _new_slide()
    _add_number(s)
    _rect(s, MARGIN, Inches(1.95), Inches(2.2), Pt(4), ACCENT, shape=MSO_SHAPE.RECTANGLE)
    tf = _box(s, MARGIN, Inches(2.1), CONTENT_W, Inches(2.2))
    rows = [
        ("Thank You", 40, True, NAVY, 6),
        ("Questions and Discussion", 18, False, ACCENT, 14),
        ("Sabinesh Rajbhandari,  Defend & Detect,  BIT Final-Year Project", 13, False, INK, 0),
    ]
    for i, (t, sz, b, c, after) in enumerate(rows):
        _para(tf, i, t, sz, b, c, after=after)


# ----------------------------------------------------------------------------
# Figure handles
# ----------------------------------------------------------------------------
F_USECASE = FIG / "fig3_1_use_case_diagram_representing_user_and_admin_bound.jpeg"
F_CLASS = FIG / "fig3_2_class_diagram_representing_service_dependencies.jpeg"
F_SEQ = FIG / "fig3_3_sequence_diagram_of_analysis_process.jpeg"
F_STATE = FIG / "fig3_4_state_diagram_showing_the_operational_lifecycle_of.jpeg"
F_ACT = FIG / "fig3_5_activity_diagram_showing_the_threat_analysis_workf.jpeg"
F_DFD = FIG / "fig3_6_dfd_showing_process_data_flow.jpeg"
F_ERD = FIG / "fig3_7_erd_of_database_schema.jpeg"
F_DASH = FIG / "fig6_1_primary_navigation_hub_and_system_dashboard.png"
F_PHISH = FIG / "fig6_2_phishing_detector_email_analysis_interface.png"
F_URL = FIG / "fig6_4_url_analyzer_link_reputation_scan_panel.png"
F_HASH = FIG / "fig6_6_file_hash_scanner_malware_hash_verification.png"
F_CVE = FIG / "fig6_8_cve_explainer_nist_nvd_integration_console.png"
F_KEV = FIG / "fig6_9_cisa_kev_active_exploitation_warning_display.png"
F_LOG = FIG / "fig6_10_security_log_translator_error_parsing_dashboard.png"

# Wipe the template's sample slides; keep its masters/layouts (footer logos).
_sldIdLst = prs.slides._sldIdLst
for sid in list(_sldIdLst):
    prs.part.drop_rel(sid.get(qn("r:id")))
    _sldIdLst.remove(sid)

title_slide()
agenda_slide()

content_slide("Introduction", "Introduction",
    "One AI workbench that unifies fragmented cyber-forensic tools.", [
    "Analysts juggle WHOIS, VirusTotal, NVD, and YARA for a single artifact — the “swivel-chair” problem.",
    "Raw outputs (hex magic-bytes, CVSS strings, server logs) are unreadable to juniors and students.",
    "Existing scanners are black boxes: a bare “Malicious / Safe” with no reasoning.",
    "Defend & Detect consolidates forensics into one slate-dark interface with plain-language reports.",
])

content_slide("Introduction", "Problem Statement",
    "Threat intelligence is fragmented, cryptic, and risky to look up.", [
    "Fragmentation: separate tools needed for basic telemetry on one artifact.",
    "Expert-knowledge gap: cryptic raw outputs need a translation layer.",
    "No pedagogy: black-box verdicts explain neither why nor how.",
    "Privacy risk: credential-leak checks expose sensitive data to online databases.",
])

content_slide("Introduction", "Objectives",
    "From scattered tools to one explainable forensic workbench.", [
    "Orchestrate real-time threat lookups across VirusTotal, NIST NVD, and CISA KEV.",
    "Classify phishing intent with a BERT deep-learning model; synthesise reports with Llama 3.",
    "Run a local YARA + binary engine (magic bytes, entropy) for deterministic file analysis.",
    "Add a privacy-safe k-anonymity breach checker; persist history in SQLite with Compare Mode.",
])

content_slide("Requirements", "Functional Requirements",
    "Eight modules: classify, look up, scan, translate, compare.", [
    "Phishing classifier (BERT + Llama 3) and URL reputation (VirusTotal, WHOIS, redirects).",
    "File-hash scanner with string extraction; CVE translator over the NIST NVD.",
    "Server-log translator; breach checker for both emails and passwords.",
    "SQLite history persistence and a side-by-side Compare Mode.",
])

content_slide("Requirements", "Non-functional Requirements",
    "Usable, fast, secure, and modular.", [
    "Usability: a mobile-responsive custom Slate dark theme.",
    "Performance: external API calls time out gracefully within fifteen seconds.",
    "Security: parameterised SQL queries and redacted password input.",
    "Modularity: a service-oriented architecture decoupling core logic from the UI.",
])

content_slide("Analysis & Design", "System Architecture",
    "A service-oriented design that keeps the UI responsive.", [
    "StreamlitApp is the UI controller; Python async keeps every lookup non-blocking.",
    "IntelligenceService runs local checks (YARA, entropy) and routes the external API queries.",
    "GroqService compiles the Llama 3 report; DatabaseService manages the SQLite connection pool.",
    "Flow: upload → local + external checks → AI synthesis → saved to SQLite → rendered.",
])

image_slide("Analysis & Design", "Use-Case Diagram", [F_USECASE])
image_slide("Analysis & Design", "Class & Sequence Diagrams", [F_CLASS, F_SEQ],
            ["Class diagram (service dependencies)", "Sequence diagram (analysis process)"])
image_slide("Analysis & Design", "State & Activity Diagrams", [F_STATE, F_ACT],
            ["State diagram (operational lifecycle)", "Activity diagram (analysis workflow)"])
image_slide("Analysis & Design", "Data Flow & Entity-Relationship", [F_DFD, F_ERD],
            ["DFD (process data flow)", "ER diagram (database schema)"])

pipeline_slide("Algorithm", "The Analysis Pipeline",
    "One artifact flows from input to an AI-written report.", [
    ("Input", "Submit a URL, hash, email, or log; format is checked."),
    ("Analyse (Fork)", "Parallel local YARA + entropy, API lookups, BERT."),
    ("Merge (Join)", "The three result paths combine into one record."),
    ("AI Synthesis", "Llama 3 turns the findings into a readable report."),
    ("Output", "Gauges + report rendered, then saved to SQLite."),
])

cards_slide("Key Terms", "Key Terms (1 of 2)",
    "The AI and detection engines behind the platform.", [
    ("BERT", "Transformer that reads an email’s intent, not just keywords."),
    ("Llama 3 (via Groq)", "LLM that writes the plain-language threat report, fast."),
    ("YARA", "Signature engine matching text and binary malware patterns."),
    ("Shannon Entropy", "Randomness score that flags DGA / suspicious domains."),
])

cards_slide("Key Terms", "Key Terms (2 of 2)",
    "The forensic and threat-intelligence concepts used.", [
    ("Magic Bytes", "Leading hex bytes revealing a file’s true type (4D 5A = PE)."),
    ("k-Anonymity", "Sends only a SHA-1 prefix to check breaches privately."),
    ("CVSS / EPSS", "Severity score and 30-day exploit-probability of a CVE."),
    ("CISA KEV / NVD", "Feeds for actively-exploited and catalogued vulnerabilities."),
])

content_slide("Implementation", "Implementation Tools",
    "An open-source, AI-first stack with no paid software.", [
    "Python 3.12 + Streamlit: the Slate dark web workbench (UI).",
    "Groq SDK (Llama 3) + Hugging Face BERT: AI reasoning and phishing classification.",
    "yara-python, python-whois, PyPDF2: the static forensic engines.",
    "VirusTotal, NIST NVD, CISA KEV, HaveIBeenPwned APIs; SQLite for persistence.",
])

content_slide("Implementation", "Module Details",
    "Eight decoupled micro-services, one per forensic task.", [
    "Phishing Detector & URL Analyzer: BERT intent, WHOIS / redirects, entropy DGA detection.",
    "Hash Scanner & CVE Explainer: VirusTotal hashes, magic bytes, NVD / CVSS, CISA KEV alerts.",
    "Log Translator & Breach Checker: LLM log parsing; k-anonymity password / email exposure.",
    "Compare Mode & History: side-by-side diffs; SQLite history with CSV / PDF export.",
])

stat_image_slide("Results", "Results: Module Testing", F_DASH,
    "Eight unit tests across all modules — 100% pass rate.", [
    ("8 / 8", "unit test cases passed"),
    ("100%", "module pass rate (supervisor log)"),
    ("70+", "VirusTotal flags on the EICAR test file"),
    ("10.0", "CVSS parsed (Log4Shell, CVE-2021-44228)"),
    ], verdict="BERT flags phishing email > 90% confidence")

image_slide("Results", "Platform in Action (1 of 2)", [F_PHISH, F_URL, F_HASH],
            ["Phishing Detector", "URL Analyzer", "File Hash Scanner"])
image_slide("Results", "Platform in Action (2 of 2)", [F_CVE, F_KEV, F_LOG],
            ["CVE Explainer (NVD)", "CISA KEV alert", "Log Translator"])

content_slide("Limitations", "Limitations",
    "Honest about what a static workbench does not do.", [
    "Static analysis and point-in-time lookups only — no dynamic sandboxing.",
    "It does not execute live malware or observe runtime behaviour.",
    "Not an active EDR agent running on the host.",
    "Throughput is bound by third-party public-API rate limits.",
])

content_slide("Conclusion", "Conclusion",
    "Deterministic rules + ML + generative AI, in one workbench.", [
    "Combines YARA / magic bytes, BERT classification, and Llama 3 reasoning.",
    "Eliminates the swivel-chair workflow across fragmented tools.",
    "A service-oriented design keeps the logic modular and the UI clean.",
    "Future work: SIEM connectors, auto-generated YARA rules, and sandbox integration.",
])

closing_slide()

prs.save(str(OUT))
print(f"Wrote {OUT}  ({len(prs.slides)} slides)")
