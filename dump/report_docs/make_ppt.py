"""Build the PocketDRS final-defense presentation.

Reproducible: re-run to regenerate pocketdrs_presentation.pptx. Uses the
institutional BIT template theme for background/footer logos, but lays out every
slide with explicit, controlled geometry so nothing clips or overlaps.

Design language: clean, minimal, modern. A small palette (navy + one accent
blue), a thin accent rule under every title, a section tag for orientation, and
generous whitespace. Text is kept short; meaning is carried by visual blocks
instead of long bullets:
  - the algorithm is a five-card pipeline strip,
  - the key terms are a two-by-two card grid (term + one plain line),
  - the results are big-number stat callouts beside the figure.
Every acronym is given its full form, and the hard computer-vision terms get two
"Key Terms" slides. Requires python-pptx + opencv.
"""

from __future__ import annotations

from pathlib import Path

import cv2
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

BASE = Path(__file__).resolve().parent
TEMPLATE = Path("/Users/nirajkafle/Desktop/niraj/dev-projects/pocket-drs/dump/bit_template.pptx")
FIG = BASE / "figures"
OUT = BASE / "pocketdrs_presentation.pptx"

# Restrained, modern palette.
NAVY = RGBColor(0x12, 0x33, 0x5B)   # titles, lead text, big stats
INK = RGBColor(0x2A, 0x2A, 0x2A)    # body text
ACCENT = RGBColor(0x2E, 0x6F, 0xB5) # rules, tags, badges, accents
MUTE = RGBColor(0x7C, 0x86, 0x90)   # captions, slide numbers, stat labels
CARD = RGBColor(0xF1, 0xF5, 0xFA)   # soft card fill
HAIR = RGBColor(0xD7, 0xE0, 0xEC)   # hairline borders / chevrons
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation(str(TEMPLATE))
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[10]

MARGIN = Inches(0.5)
TAG_TOP = Inches(0.42)
TAG_H = Inches(0.22)
TITLE_TOP = Inches(0.62)
TITLE_H = Inches(0.52)
RULE_TOP = Inches(1.16)
BODY_TOP = Inches(1.40)
FOOTER_RESERVE = Inches(0.86)
BODY_H = SH - BODY_TOP - FOOTER_RESERVE
CONTENT_W = SW - 2 * MARGIN

_slide_no = 0


def _new_slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    return s


def _para(tf, i, text, size, bold, color, after=0.0, align=PP_ALIGN.LEFT, lh=1.04):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = text
    p.alignment = align
    p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color
    p.space_after = Pt(after); p.space_before = Pt(0); p.line_spacing = lh
    return p


def _box(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    return tf


def _rect(slide, x, y, w, h, fill, line=None, line_w=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = line_w
    try:  # softer corners on rounded rectangles
        sp.adjustments[0] = 0.08
    except (IndexError, KeyError):
        pass
    return sp


def _add_number(slide) -> None:
    global _slide_no
    _slide_no += 1
    if _slide_no == 1:
        return
    tf = _box(slide, SW - MARGIN - Inches(0.7), SH - Inches(0.4), Inches(0.7), Inches(0.3))
    _para(tf, 0, str(_slide_no), 9, False, MUTE, align=PP_ALIGN.RIGHT)


def _add_header(slide, title: str, tag: str | None) -> None:
    # Drop the section tag when it would just echo a word already in the title
    # (e.g. tag "CONCLUSION" above the title "Conclusion").
    if tag:
        title_low = title.lower()
        tag_words = [w for w in tag.lower().replace("&", " ").split() if len(w) > 3]
        if any(w in title_low for w in tag_words):
            tag = None
    if tag:
        tf = _box(slide, MARGIN, TAG_TOP, CONTENT_W, TAG_H)
        _para(tf, 0, tag.upper(), 11, True, ACCENT)
    tf = _box(slide, MARGIN, TITLE_TOP, CONTENT_W, TITLE_H, anchor=MSO_ANCHOR.MIDDLE)
    _para(tf, 0, title, 23, True, NAVY)
    _rect(slide, MARGIN, RULE_TOP, Inches(1.9), Pt(3), ACCENT, shape=MSO_SHAPE.RECTANGLE)
    _add_number(slide)


def _lead(slide, text: str):
    """The single bold message of the slide; returns its bottom y."""
    h = Inches(0.5)
    tf = _box(slide, MARGIN, BODY_TOP, CONTENT_W, h, anchor=MSO_ANCHOR.MIDDLE)
    _para(tf, 0, text, 16, True, NAVY, lh=1.02)
    return BODY_TOP + h + Inches(0.12)


# ----------------------------------------------------------------------------
# Slide templates
# ----------------------------------------------------------------------------

def content_slide(tag: str, title: str, lead: str, bullets: list[str]) -> None:
    """Title + bold lead + a few short bullets with an accent marker."""
    s = _new_slide()
    _add_header(s, title, tag)
    top = _lead(s, lead)
    tf = _box(s, MARGIN, top, CONTENT_W, (BODY_TOP + BODY_H) - top)
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r1 = p.add_run(); r1.text = "▪  "
        r1.font.size = Pt(13); r1.font.color.rgb = ACCENT; r1.font.bold = True
        r2 = p.add_run(); r2.text = text
        r2.font.size = Pt(15); r2.font.color.rgb = INK
        p.space_after = Pt(11); p.line_spacing = 1.06


def pipeline_slide(tag: str, title: str, lead: str, steps: list[tuple[str, str]]) -> None:
    """Five numbered cards in a row, with chevrons, reading as a flow."""
    s = _new_slide()
    _add_header(s, title, tag)
    top = _lead(s, lead)
    n = len(steps)
    gap = Inches(0.16)
    card_w = int((CONTENT_W - gap * (n - 1)) / n)
    band_top = top + Inches(0.08)
    band_h = int((BODY_TOP + BODY_H) - band_top)
    card_h = min(band_h, int(Inches(2.15)))
    card_top = band_top + (band_h - card_h) // 2
    badge = Inches(0.42)
    x = MARGIN
    for k, (name, desc) in enumerate(steps):
        _rect(s, x, card_top, card_w, card_h, CARD, line=HAIR, line_w=Pt(0.75))
        # number badge
        bx = x + (card_w - badge) // 2
        _rect(s, bx, card_top + Inches(0.22), badge, badge, ACCENT, shape=MSO_SHAPE.OVAL)
        btf = _box(s, bx, card_top + Inches(0.22), badge, badge, anchor=MSO_ANCHOR.MIDDLE)
        _para(btf, 0, str(k + 1), 16, True, WHITE, align=PP_ALIGN.CENTER)
        # name + description, centred in the space below the badge
        ty = card_top + Inches(0.74)
        ttf = _box(s, x + Inches(0.1), ty, card_w - Inches(0.2),
                   (card_top + card_h) - ty - Inches(0.08), anchor=MSO_ANCHOR.MIDDLE)
        _para(ttf, 0, name, 12.5, True, NAVY, after=4, align=PP_ALIGN.CENTER, lh=1.0)
        _para(ttf, 1, desc, 10.5, False, INK, align=PP_ALIGN.CENTER, lh=1.06)
        # chevron between cards
        if k < n - 1:
            cf = _box(s, x + card_w - Inches(0.02), card_top, gap + Inches(0.04), card_h,
                      anchor=MSO_ANCHOR.MIDDLE)
            _para(cf, 0, "›", 20, True, ACCENT, align=PP_ALIGN.CENTER)
        x += card_w + gap


def cards_slide(tag: str, title: str, lead: str, cards: list[tuple[str, str]]) -> None:
    """Two-by-two grid of term cards: bold term + one plain line."""
    s = _new_slide()
    _add_header(s, title, tag)
    top = _lead(s, lead)
    cols, gap = 2, Inches(0.22)
    rows = (len(cards) + cols - 1) // cols
    cw = int((CONTENT_W - gap * (cols - 1)) / cols)
    grid_top = top + Inches(0.06)
    ch = int(((BODY_TOP + BODY_H) - grid_top - gap * (rows - 1)) / rows)
    for idx, (term, meaning) in enumerate(cards):
        r, c = divmod(idx, cols)
        x = MARGIN + c * (cw + gap)
        y = grid_top + r * (ch + gap)
        _rect(s, x, y, cw, ch, CARD, line=HAIR, line_w=Pt(0.75))
        tf = _box(s, x + Inches(0.22), y, cw - Inches(0.44), ch, anchor=MSO_ANCHOR.MIDDLE)
        _para(tf, 0, term, 14, True, NAVY, after=4, lh=1.0)
        _para(tf, 1, meaning, 12, False, INK, lh=1.07)


def _fit(img: Path, max_w: int, max_h: int) -> tuple[int, int]:
    h, w = cv2.imread(str(img)).shape[:2]
    ar = w / h
    width, height = max_w, int(max_w / ar)
    if height > max_h:
        height, width = max_h, int(max_h * ar)
    return width, height


def image_slide(tag: str, title: str, images: list[Path],
                captions: list[str] | None = None) -> None:
    images = [im for im in images if im.exists()]
    if not images:
        return
    s = _new_slide()
    _add_header(s, title, tag)
    n = len(images)
    gap = Inches(0.3)
    cap_h = Inches(0.3) if captions else Inches(0.0)
    cell_w = int((CONTENT_W - gap * (n - 1)) / n)
    area_h = int(BODY_H - cap_h)
    x = MARGIN
    for k, img in enumerate(images):
        w, h = _fit(img, cell_w, area_h)
        s.shapes.add_picture(str(img), x + (cell_w - w) // 2,
                             BODY_TOP + (area_h - h) // 2, width=w, height=h)
        if captions:
            cf = _box(s, x, BODY_TOP + area_h, cell_w, cap_h)
            _para(cf, 0, captions[k], 11, False, MUTE, align=PP_ALIGN.CENTER)
        x += cell_w + gap


def stat_image_slide(tag: str, title: str, image: Path, lead: str,
                     stats: list[tuple[str, str]], verdict: str | None = None) -> None:
    """Figure on the left; big-number stat cards (and an optional verdict pill) on the right."""
    s = _new_slide()
    _add_header(s, title, tag)
    top = _lead(s, lead)
    area_h = (BODY_TOP + BODY_H) - top
    half = int((CONTENT_W - Inches(0.35)) / 2)
    if image.exists():
        w, h = _fit(image, half, int(area_h))
        s.shapes.add_picture(str(image), MARGIN + (half - w) // 2,
                             top + (int(area_h) - h) // 2, width=w, height=h)
    # right column
    rx = MARGIN + half + Inches(0.35)
    gap = Inches(0.16)
    pill_h = Inches(0.62) if verdict else Inches(0.0)
    pill_gap = Inches(0.16) if verdict else Inches(0.0)
    grid_h = int(area_h - pill_h - pill_gap)
    rows = (len(stats) + 1) // 2
    cw = int((half - gap) / 2)
    chh = int((grid_h - gap * (rows - 1)) / rows)
    for idx, (val, label) in enumerate(stats):
        r, c = divmod(idx, 2)
        x = rx + c * (cw + gap)
        y = top + r * (chh + gap)
        _rect(s, x, y, cw, chh, CARD, line=HAIR, line_w=Pt(0.75))
        tf = _box(s, x + Inches(0.12), y, cw - Inches(0.24), chh, anchor=MSO_ANCHOR.MIDDLE)
        _para(tf, 0, val, 21, True, NAVY, after=2, lh=0.95)
        _para(tf, 1, label, 10.5, False, MUTE, lh=1.0)
    if verdict:
        py = top + area_h - pill_h
        _rect(s, rx, py, half, pill_h, ACCENT)
        tf = _box(s, rx + Inches(0.2), py, half - Inches(0.4), pill_h, anchor=MSO_ANCHOR.MIDDLE)
        _para(tf, 0, verdict, 14, True, WHITE, align=PP_ALIGN.CENTER, lh=1.0)


def hero_image_slide(tag: str, title: str, image: Path, hero: str, hero_sub: str,
                     points: list[str]) -> None:
    """Figure on the left; one giant headline stat and a few short points on the right."""
    s = _new_slide()
    _add_header(s, title, tag)
    area_h = BODY_H
    half = int((CONTENT_W - Inches(0.35)) / 2)
    if image.exists():
        w, h = _fit(image, half, int(area_h))
        s.shapes.add_picture(str(image), MARGIN + (half - w) // 2,
                             BODY_TOP + (int(area_h) - h) // 2, width=w, height=h)
    rx = MARGIN + half + Inches(0.35)
    tf = _box(s, rx, BODY_TOP, half, Inches(1.5))
    _para(tf, 0, hero, 54, True, ACCENT, after=0, lh=0.9)
    _para(tf, 1, hero_sub, 14, True, NAVY, lh=1.0)
    tf2 = _box(s, rx, BODY_TOP + Inches(1.55), half, area_h - Inches(1.55))
    for i, text in enumerate(points):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        r1 = p.add_run(); r1.text = "▪  "
        r1.font.size = Pt(12); r1.font.color.rgb = ACCENT; r1.font.bold = True
        r2 = p.add_run(); r2.text = text
        r2.font.size = Pt(13); r2.font.color.rgb = INK
        p.space_after = Pt(9); p.line_spacing = 1.06


def title_slide() -> None:
    s = _new_slide()
    _add_number(s)
    _rect(s, MARGIN, Inches(1.05), Inches(2.2), Pt(4), ACCENT, shape=MSO_SHAPE.RECTANGLE)
    tf = _box(s, MARGIN, Inches(1.2), CONTENT_W, Inches(3.2))
    rows = [
        ("PocketDRS", 42, True, NAVY, 6),
        ("A Single-View 3D Trajectory Reconstruction and", 17, False, INK, 0),
        ("Decision Review System (DRS) for Cricket", 17, False, INK, 18),
        ("Niraj Kafle,  BIT 7th Semester,  ID LC0003001674", 14, False, INK, 3),
        ("Supervisor: Mr. Saishab Bhattarai", 13, False, INK, 3),
        ("Phoenix College of Management,  Lincoln University College", 13, False, ACCENT, 0),
    ]
    for i, (t, sz, b, c, after) in enumerate(rows):
        _para(tf, i, t, sz, b, c, after=after, lh=1.02)


def agenda_slide() -> None:
    s = _new_slide()
    _add_header(s, "Agenda", None)
    items = [
        "Introduction and the Problem",
        "Objectives and Requirements",
        "System Analysis and Design",
        "The Five-Step Algorithm",
        "Key Terms",
        "Implementation and Tools",
        "Results and Validation",
        "Limitations and Conclusion",
    ]
    tf = _box(s, MARGIN, BODY_TOP, CONTENT_W, BODY_H, anchor=MSO_ANCHOR.MIDDLE)
    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run(); run.text = f"{i + 1:02d}    "
        run.font.size = Pt(16); run.font.bold = True; run.font.color.rgb = ACCENT
        run2 = p.add_run(); run2.text = text
        run2.font.size = Pt(16); run2.font.color.rgb = INK
        p.space_after = Pt(8); p.line_spacing = 1.05


def closing_slide() -> None:
    s = _new_slide()
    _add_number(s)
    _rect(s, MARGIN, Inches(1.95), Inches(2.2), Pt(4), ACCENT, shape=MSO_SHAPE.RECTANGLE)
    tf = _box(s, MARGIN, Inches(2.1), CONTENT_W, Inches(2.2))
    rows = [
        ("Thank You", 40, True, NAVY, 6),
        ("Questions and Discussion", 18, False, ACCENT, 14),
        ("Niraj Kafle,  PocketDRS,  BIT Final-Year Project", 13, False, INK, 0),
    ]
    for i, (t, sz, b, c, after) in enumerate(rows):
        _para(tf, i, t, sz, b, c, after=after)


# Wipe the template's sample slides; keep its masters/layouts (footer logos).
_sldIdLst = prs.slides._sldIdLst
for sid in list(_sldIdLst):
    prs.part.drop_rel(sid.get(qn("r:id")))
    _sldIdLst.remove(sid)

title_slide()
agenda_slide()

content_slide("Introduction", "Introduction",
    "Ball-tracking review for everyday cricket, from a single phone.", [
    "Leg Before Wicket (LBW): cricket's hardest call, judged live from one angle.",
    "Hawk-Eye needs six to eight synced high-speed cameras and costly hardware.",
    "Out of reach for schools, clubs, and training grounds.",
    "PocketDRS rebuilds the delivery in 3D from one phone clip and supports the call.",
])

content_slide("Introduction", "Problem Statement",
    "Recover a true 3D ball path from a single flat (2D) video.", [
    "An umpire judges line, length, height, and impact at once, in real time.",
    "Multi-camera rigs are accurate but unaffordable for grassroots cricket.",
    "One phone gives flat 2D only: unknown depth, motion blur, occlusion.",
    "Goal: a 3D flight and a clear, explained verdict from one viewpoint.",
])

content_slide("Introduction", "Objectives",
    "A phone-only pipeline: from clip to verdict.", [
    "Recover the camera position from the marked stumps (Perspective-n-Point, PnP).",
    "Detect and track the ball (motion + colour + a trained YOLO detector).",
    "Link the detections into one smooth path (Random Sample Consensus, RANSAC).",
    "Rebuild the 3D path with projectile physics; apply the three ICC Rule 36 checks.",
])

content_slide("Requirements", "Functional Requirements",
    "Twelve requirements: capture, calibrate, analyse, visualise.", [
    "Record or upload a clip and tap the four pitch corners to calibrate.",
    "Detect the ball per frame and link it into one trajectory.",
    "Recover the camera position and rebuild the path in real-world 3D.",
    "Return an LBW verdict with a 2D overlay and an interactive 3D view.",
])

content_slide("Requirements", "Non-functional Requirements",
    "Fast, easy, clear, and portable, with no special hardware.", [
    "Speed: a delivery analysed in about thirty seconds.",
    "Ease of use: four taps, no setup, no special hardware.",
    "Transparency: every verdict shows all three checks and the predicted impact.",
    "Portability: a normal Android or iPhone plus a small cloud server.",
])

image_slide("Analysis & Design", "Use-Case Diagram", [FIG / "use_case_diagram.png"])
image_slide("Analysis & Design", "Data Model: Entity-Relationship (ER) Diagram", [FIG / "er_diagram.png"])
image_slide("Analysis & Design", "Process Model: Data Flow Diagrams (DFD)",
            [FIG / "dfd_level0.png", FIG / "dfd_level1.png"],
            ["Level 0 (context)", "Level 1"])
image_slide("Analysis & Design", "System Architecture", [FIG / "architecture.png"])

pipeline_slide("Algorithm", "Algorithm: The Five-Step Pipeline",
    "Five steps turn one phone clip into an LBW verdict.", [
    ("Calibrate", "Marked stumps reveal the camera position (PnP)."),
    ("Detect", "Find the ball each frame: motion, colour, YOLO."),
    ("Link", "RANSAC fits one curve, drops false detections."),
    ("Reconstruct", "Gravity-based fit lifts the track to a 3D arc."),
    ("Decide", "Three ICC Rule 36 checks give the verdict."),
])

cards_slide("Key Terms", "Key Terms (1 of 2)",
    "The standard computer-vision tools behind the pipeline.", [
    ("Camera pose (PnP)", "Finds where the phone stood, from the marked stumps."),
    ("Homography", "Maps any flat-pitch point in the photo to real metres."),
    ("RANSAC", "Fits the ball's path while ignoring wrong detections."),
    ("Projectile motion", "Gravity's curve; forces a real ball flight, not noise."),
])

cards_slide("Key Terms", "Key Terms (2 of 2)",
    "How the ball is found and depth recovered from one view.", [
    ("MOG2", "Separates the moving ball from the still background."),
    ("HSV colour", "A colour space that isolates one ball colour despite light."),
    ("YOLO", "A trained network that spots the cricket ball in a frame."),
    ("Depth-from-size", "Smaller ball looks farther; recovers 3D from one camera."),
])

content_slide("Implementation", "Implementation Tools",
    "An open-source phone-to-cloud stack, with no paid software.", [
    "Flutter + Dart: the cross-platform mobile app.",
    "FastAPI (Python): the analysis server, over a REST web interface.",
    "OpenCV, NumPy, SciPy: the computer-vision maths.",
    "YOLO + PyTorch: the trained ball detector. Three.js: the 3D viewer. Firebase: sign-in and storage.",
])

content_slide("Implementation", "Module Details",
    "One module per pipeline step, each small and testable.", [
    "Calibration: camera position from the tapped stumps; also gives pitch length.",
    "Detection: motion + colour + YOLO inside the pitch area.",
    "Trajectory: RANSAC seed, refine, then merge the pre- and post-bounce arcs.",
    "Reconstruction + decision: lift to 3D, predict to the stumps, apply ICC Rule 36.",
])

stat_image_slide("Results", "Real-Video Test: test3.mp4", FIG / "test3_overlay.png",
    "A real handheld off-spin clip, tracked end to end (full 20.12 m pitch).", [
    ("29", "ball detections (28 fit the path)"),
    ("10.55 px", "calibration error"),
    ("84.0 km/h", "release speed"),
    ("2.3°", "off-break turn (spin)"),
    ], verdict="NOT OUT  ·  predicted to miss the stumps")

image_slide("Results", "Interactive 3D Ball-Path View: test3", [FIG / "test3_3d_path.png"])

hero_image_slide("Results", "Synthetic Validation: 8 Scenarios", FIG / "synth_summary.png",
    "8/8", "scenarios passed", [
    "Controlled deliveries: medium pace, middle-stump line, full length.",
    "Camera position recovered correctly in every scene.",
    "All within honest single-camera limits (speed, bounce, impact).",
    "Every delivery was truly OUT; the system returned OUT for all eight.",
])

content_slide("Limitations", "Limitations",
    "Honest about what one camera can and cannot do.", [
    "One phone cannot triangulate depth like six to eight synced Hawk-Eye cameras.",
    "Depth from ball size is good to a few tens of centimetres, not millimetres.",
    "The 8/8 passes are within single-camera limits, not broadcast limits.",
    "Built as decision support for nets and clubs, not an official broadcast replacement.",
])

content_slide("Conclusion", "Conclusion",
    "One phone, a known pitch, and simple physics: enough for usable LBW support.", [
    "A full pipeline: calibrate, detect, track, reconstruct, decide.",
    "Validated on synthetic deliveries (8/8) and real footage.",
    "Refuses to invent a verdict on unusable clips: it fails clearly, not silently.",
    "Future work: two-phone capture, learned bounce model, and on-phone detection.",
])

closing_slide()

prs.save(str(OUT))
print(f"Wrote {OUT}  ({len(prs.slides)} slides)")
